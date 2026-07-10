#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
跨平台文件内容字典检索工具 - 图形界面版 V5.1
支持：图片OCR、扫描件PDF OCR、OFD文字提取
修复：OCR中文空格、进度条实时更新
"""

import os, sys, re, csv, time, zipfile, tempfile, shutil, threading
import subprocess, platform, struct, io
from datetime import datetime
from pathlib import Path
from html import escape

# ==================== GUI ====================
import tkinter as tk
from tkinter import ttk, filedialog, messagebox

# ==================== 资源路径 ====================
def get_resource_path(relative_path):
    try:
        base_path = sys._MEIPASS
    except AttributeError:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)

# ==================== 配置常量 ====================
TEXT_EXTENSIONS = {
    '.txt', '.csv', '.log', '.xml', '.json', '.html', '.htm',
    '.md', '.py', '.js', '.java', '.c', '.cpp', '.h', '.hpp',
    '.css', '.sql', '.yaml', '.yml', '.ini', '.cfg', '.conf',
    '.sh', '.bat', '.ps1', '.rtf', '.tex', '.rst', '.php',
    '.asp', '.aspx', '.jsp', '.toml', '.properties'
}
DOCX_EXTENSIONS = {'.docx'}
XLSX_EXTENSIONS = {'.xlsx', '.xlsm'}
PPTX_EXTENSIONS = {'.pptx', '.pptm'}
DOC_EXTENSIONS = {'.doc'}
XLS_EXTENSIONS = {'.xls'}
PPT_EXTENSIONS = {'.ppt'}
PDF_EXTENSIONS = {'.pdf'}
ZIP_EXTENSIONS = {'.zip'}
EML_EXTENSIONS = {'.eml'}
WPS_TEXT_EXTENSIONS = {'.wps', '.wpt'}
WPS_SHEET_EXTENSIONS = {'.et', '.ett'}
WPS_SLIDE_EXTENSIONS = {'.dps', '.dpt'}
IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.bmp', '.tif', '.tiff'}
OFD_EXTENSIONS = {'.ofd'}

SKIP_DIRS = {
    'System Volume Information', '$RECYCLE.BIN', 'RECYCLER',
    'Windows', 'Program Files', 'Program Files (x86)',
    'ProgramData', 'AppData', '.git', '.svn', '__pycache__',
    'node_modules', '.cache', 'proc', 'sys', 'dev', 'run',
    '/proc', '/sys', '/dev', '/run', '/boot', '/etc/ssl'
}

SKIP_EXTENSIONS = {
    '.exe', '.dll', '.so', '.dylib', '.bin', '.dat', '.db',
    '.mdb', '.accdb', '.pst', '.ost', '.iso', '.img', '.vmdk',
    '.vhd', '.vdi', '.gho', '.wim', '.tar', '.gz', '.bz2',
    '.7z', '.rar', '.mp3', '.mp4', '.avi', '.mkv', '.wmv',
    '.mov', '.flv', '.ttf', '.otf', '.woff', '.woff2', '.eot',
    '.pyc', '.pyo', '.class', '.o', '.obj', '.lib', '.a',
    '.sys', '.drv', '.ocx', '.msi', '.cab'
}
MAX_FILE_SIZE_DEFAULT = 50 * 1024 * 1024

# ==================== 文本提取模块 ====================
class TextExtractor:

    @staticmethod
    def extract_text(filepath, ext_lower):
        if ext_lower in TEXT_EXTENSIONS:
            return TextExtractor._extract_text_file(filepath)
        if ext_lower in DOCX_EXTENSIONS:
            return TextExtractor._extract_docx(filepath)
        if ext_lower in XLSX_EXTENSIONS:
            return TextExtractor._extract_xlsx(filepath)
        if ext_lower in PPTX_EXTENSIONS:
            return TextExtractor._extract_pptx(filepath)
        if ext_lower in PDF_EXTENSIONS:
            return TextExtractor._extract_pdf(filepath)
        if ext_lower in EML_EXTENSIONS:
            return TextExtractor._extract_text_file(filepath)
        if ext_lower in DOC_EXTENSIONS:
            return TextExtractor.extract_text_doc(filepath)
        if ext_lower in (XLS_EXTENSIONS | PPT_EXTENSIONS |
                         WPS_TEXT_EXTENSIONS | WPS_SHEET_EXTENSIONS | WPS_SLIDE_EXTENSIONS):
            return TextExtractor.extract_text_legacy(filepath, ext_lower)
        if ext_lower in IMAGE_EXTENSIONS:
            return TextExtractor._extract_image_ocr(filepath)
        if ext_lower in OFD_EXTENSIONS:
            return TextExtractor._extract_ofd(filepath)
        return "", f"不支持的文件格式: {ext_lower}"

    @staticmethod
    def _extract_text_file(filepath):
        encodings = ['utf-8', 'gbk', 'gb2312', 'gb18030', 'latin-1', 'cp1252']
        for enc in encodings:
            try:
                with open(filepath, 'r', encoding=enc, errors='strict') as f:
                    return f.read(), None
            except (UnicodeDecodeError, UnicodeError):
                continue
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read(), None
        except Exception as e:
            return "", f"无法读取文件: {e}"

    @staticmethod
    def _extract_docx(filepath):
        try:
            from docx import Document
            doc = Document(filepath)
            paragraphs = [p.text for p in doc.paragraphs]
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        paragraphs.append(cell.text)
            return '\n'.join(paragraphs), None
        except ImportError:
            return "", "需要 python-docx"
        except Exception as e:
            return "", f"读取docx失败: {e}"

    @staticmethod
    def _extract_xlsx(filepath):
        try:
            from openpyxl import load_workbook
            wb = load_workbook(filepath, read_only=True, data_only=True)
            texts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                for row in ws.iter_rows(values_only=True):
                    row_text = ' '.join([str(c) if c is not None else '' for c in row])
                    if row_text.strip():
                        texts.append(row_text)
            wb.close()
            return '\n'.join(texts), None
        except ImportError:
            return "", "需要 openpyxl"
        except Exception as e:
            return "", f"读取xlsx失败: {e}"

    @staticmethod
    def _extract_pptx(filepath):
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            texts = [shape.text for slide in prs.slides
                     for shape in slide.shapes if hasattr(shape, 'text') and shape.text]
            return '\n'.join(texts), None
        except ImportError:
            return "", "需要 python-pptx"
        except Exception as e:
            return "", f"读取pptx失败: {e}"

    @staticmethod
    def _extract_pdf(filepath):
        text = ""
        try:
            import pdfplumber
            with pdfplumber.open(filepath) as pdf:
                texts = [page.extract_text() for page in pdf.pages if page.extract_text()]
            text = '\n'.join(texts)
        except:
            pass
        if not text.strip():
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(filepath)
                texts = [page.extract_text() for page in reader.pages if page.extract_text()]
                text = '\n'.join(texts)
            except:
                pass
        if text.strip() and len(text.strip()) > 20:
            return text, None
        else:
            return TextExtractor._extract_pdf_ocr(filepath)

    @staticmethod
    def _extract_pdf_ocr(filepath):
        try:
            from pdf2image import convert_from_path
            # 150 DPI + 只扫前 5 页
            images = convert_from_path(filepath, dpi=150, first_page=1, last_page=10)
            texts = []
            for i, img in enumerate(images):
                page_text = TextExtractor._ocr_image(img)
                if page_text.strip():
                    texts.append(page_text)
            result = '\n'.join(texts)
            return (result, None) if result.strip() else ("", "PDF OCR 未识别到文字")
        except ImportError:
            return "", "需要 pdf2image"
        except Exception as e:
            return "", f"PDF OCR 失败: {e}"

    @staticmethod
    def _extract_image_ocr(filepath):
        try:
            from PIL import Image
            img = Image.open(filepath)
            text = TextExtractor._ocr_image(img)
            return (text, None) if text.strip() else ("", "图片未识别到文字")
        except ImportError:
            return "", "需要 Pillow: pip install Pillow"
        except Exception as e:
            return "", f"图片读取失败: {e}"

    @staticmethod
    def _ocr_image(pil_image):
        tesseract_exe, tessdata_dir = TextExtractor._get_tesseract_path()
        if not tesseract_exe:
            return ""
        try:
            import pytesseract
            pytesseract.pytesseract.tesseract_cmd = tesseract_exe
            if tessdata_dir:
                os.environ['TESSDATA_PREFIX'] = tessdata_dir
            # 改为 psm 4（单列可变大小文字，速度更快）
            text = pytesseract.image_to_string(pil_image, lang='chi_sim+eng', config='--psm 4')
            # 去除 CJK 字符之间的多余空格
            text = re.sub(r'(?<=[\u4e00-\u9fff\u3400-\u4dbf])\s+(?=[\u4e00-\u9fff\u3400-\u4dbf])', '', text)
            return text
        except ImportError:
            return ""
        except Exception as e:
            return f"[OCR 错误: {e}]"

    @staticmethod
    def _get_tesseract_path():
        try:
            base = sys._MEIPASS
        except AttributeError:
            base = os.path.dirname(os.path.abspath(__file__))
        candidates = []
        if platform.system() == 'Windows':
            candidates = [
                os.path.join(base, 'tesseract', 'tesseract.exe'),
                os.path.join(base, 'tesseract.exe'),
                r'C:\Program Files\Tesseract-OCR\tesseract.exe',
                r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
            ]
        else:
            candidates = [
                os.path.join(base, 'tesseract', 'tesseract'),
                '/usr/bin/tesseract',
                '/usr/local/bin/tesseract',
            ]
        exe = None
        for c in candidates:
            if os.path.exists(c) and (platform.system() != 'Linux' or os.access(c, os.X_OK)):
                exe = c
                break
        if not exe:
            for cmd in ['tesseract', 'tesseract.exe']:
                if shutil.which(cmd):
                    exe = shutil.which(cmd)
                    break
        if not exe:
            return None, None

        tessdata_dir = None
        exe_dir = os.path.dirname(os.path.abspath(exe))
        test_dirs = [
            os.path.join(exe_dir, 'tessdata'),
            os.path.join(base, 'tesseract', 'tessdata'),
            os.path.join(base, 'tessdata'),
            '/usr/share/tesseract-ocr/4.00/tessdata',
            '/usr/share/tesseract-ocr/tessdata',
            '/usr/share/tessdata',
        ]
        for d in test_dirs:
            if os.path.exists(os.path.join(d, 'chi_sim.traineddata')):
                tessdata_dir = d
                break
        if not tessdata_dir:
            env_prefix = os.environ.get('TESSDATA_PREFIX', '')
            if env_prefix and os.path.exists(os.path.join(env_prefix, 'chi_sim.traineddata')):
                tessdata_dir = env_prefix

        return exe, tessdata_dir

    # =================== OFD 支持 ===================
    @staticmethod
    def _extract_ofd(filepath):
        text, err = TextExtractor._extract_ofd_text(filepath)
        if text.strip() and len(text.strip()) > 20:
            return text, None
        text2, err2 = TextExtractor._extract_ofd_ocr(filepath)
        if text2.strip():
            return text2, None
        return "", f"OFD 提取失败: 文字={err}, OCR={err2}"

    @staticmethod
    def _extract_ofd_text(filepath):
        try:
            text_parts = []
            with zipfile.ZipFile(filepath, 'r') as zf:
                for name in zf.namelist():
                    if name.lower().endswith('content.xml'):
                        xml_content = zf.read(name).decode('utf-8')
                        codes = re.findall(r'<ofd:TextCode[^>]*>(.*?)</ofd:TextCode>',
                                           xml_content, re.DOTALL)
                        text_parts.extend(codes)
            result = ''.join(text_parts)
            if result.strip():
                return result, None
            return "", "OFD 中未找到 TextObject"
        except Exception as e:
            return "", f"OFD XML 解析失败: {e}"

    @staticmethod
    def _extract_ofd_ocr(filepath):
        try:
            from PIL import Image
            texts = []
            with zipfile.ZipFile(filepath, 'r') as zf:
                image_exts = {'.jpg', '.jpeg', '.png', '.bmp', '.tif', '.tiff', '.gif'}
                for name in zf.namelist():
                    ext = os.path.splitext(name)[1].lower()
                    if ext in image_exts:
                        try:
                            data = zf.read(name)
                            img = Image.open(io.BytesIO(data))
                            img_text = TextExtractor._ocr_image(img)
                            if img_text.strip():
                                texts.append(img_text)
                        except:
                            pass
            result = '\n'.join(texts)
            return (result, None) if result.strip() else ("", "OFD 内嵌图片 OCR 未识别到文字")
        except ImportError:
            return "", "需要 Pillow"
        except Exception as e:
            return "", f"OFD OCR 失败: {e}"

    # =================== .doc 支持 ===================
    @staticmethod
    def _get_antiword_path():
        try:
            base = sys._MEIPASS
        except AttributeError:
            base = os.path.dirname(os.path.abspath(__file__))
        if platform.system() == 'Windows':
            exe_candidates = [
                os.path.join(base, 'antiword_bin', 'antiword.exe'),
                os.path.join(base, 'antiword.exe'),
            ]
        else:
            exe_candidates = [
                os.path.join(base, 'antiword_bin', 'antiword'),
                os.path.join(base, 'antiword'),
            ]
        antiword_exe = None
        for c in exe_candidates:
            if os.path.exists(c) and os.access(c, os.X_OK):
                antiword_exe = c
                break
        if not antiword_exe:
            for cmd in ['antiword', 'antiword.exe']:
                if shutil.which(cmd):
                    antiword_exe = cmd
                    break
        if not antiword_exe:
            return None, None
        mapping_dir = os.path.dirname(os.path.abspath(antiword_exe))
        if not os.path.exists(os.path.join(mapping_dir, '8859-1.txt')):
            for subdir in [os.path.join(base, 'antiword_bin'),
                           os.path.join(base),
                           '/usr/share/antiword']:
                if os.path.exists(os.path.join(subdir, '8859-1.txt')):
                    mapping_dir = subdir
                    break
        return antiword_exe, mapping_dir

    @staticmethod
    def _extract_doc_via_antiword(filepath):
        antiword_exe, mapping_dir = TextExtractor._get_antiword_path()
        if not antiword_exe:
            return "", "antiword 不可用"
        try:
            cmd = [antiword_exe, '-m', mapping_dir, '-w', '0', filepath]
            result = subprocess.run(cmd, capture_output=True, text=True,
                                    timeout=15, encoding='utf-8', errors='replace')
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout, None
            else:
                return "", f"antiword 错误: {result.stderr.strip() or '无文本输出'}"
        except subprocess.TimeoutExpired:
            return "", "antiword 转换超时"
        except FileNotFoundError:
            return "", "antiword 可执行文件未找到"
        except Exception as e:
            return "", f"antiword 异常: {e}"

    @staticmethod
    def _extract_doc_via_olefile(filepath):
        try:
            import olefile
            ole = olefile.OleFileIO(filepath)
            text_parts = []
            for stream_name in ['WordDocument', '1Table', '0Table']:
                if ole.exists(stream_name):
                    data = ole.openstream(stream_name).read()
                    extracted = TextExtractor._extract_readable_from_binary(data)
                    if extracted.strip():
                        text_parts.append(extracted)
            ole.close()
            result = '\n'.join(text_parts)
            return (result, None) if result.strip() else ("", "olefile 未提取到有效文本")
        except ImportError:
            return "", "需要 olefile"
        except Exception as e:
            return "", f"olefile 异常: {e}"

    @staticmethod
    def _extract_readable_from_binary(data):
        result_chars = []
        for i in range(0, len(data) - 1, 2):
            try:
                char_code = struct.unpack_from('<H', data, i)[0]
            except struct.error:
                break
            if 0x20 <= char_code <= 0x7E:
                result_chars.append(chr(char_code))
            elif 0x4E00 <= char_code <= 0x9FFF:
                result_chars.append(chr(char_code))
            elif 0x3400 <= char_code <= 0x4DBF:
                result_chars.append(chr(char_code))
            elif 0x3000 <= char_code <= 0x303F:
                result_chars.append(chr(char_code))
            elif 0xFF00 <= char_code <= 0xFFEF:
                result_chars.append(chr(char_code))
            elif char_code == 0x0D:
                result_chars.append('\r')
            elif char_code == 0x0A:
                result_chars.append('\n')
            elif char_code == 0x09:
                result_chars.append('\t')
            else:
                if result_chars and result_chars[-1] != ' ':
                    result_chars.append(' ')
        raw = ''.join(result_chars)
        raw = re.sub(r' {2,}', ' ', raw)
        raw = re.sub(r'\n\s*\n', '\n\n', raw)
        return raw.strip()

    @staticmethod
    def extract_text_doc(filepath):
        text, err = TextExtractor._extract_doc_via_antiword(filepath)
        if text.strip():
            return text, None
        text2, err2 = TextExtractor._extract_doc_via_olefile(filepath)
        if text2.strip():
            return text2, None
        return "", f"antiword: {err}; olefile: {err2}"

    # =================== 旧版/WPS ===================
    @staticmethod
    def _convert_via_libreoffice(filepath, outdir, to_format='txt'):
        try:
            cmd = ['soffice', '--headless', '--convert-to', to_format, '--outdir', outdir, filepath]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            basename = os.path.splitext(os.path.basename(filepath))[0]
            out_ext = '.' + to_format.split(':')[-1] if ':' in to_format else '.' + to_format
            out_path = os.path.join(outdir, basename + out_ext)
            return out_path if os.path.exists(out_path) else None
        except:
            return None

    @staticmethod
    def extract_text_legacy(filepath, ext_lower):
        if ext_lower in XLS_EXTENSIONS:
            try:
                import xlrd
                workbook = xlrd.open_workbook(filepath)
                texts = []
                for sheet in workbook.sheets():
                    for row_idx in range(sheet.nrows):
                        row_values = sheet.row_values(row_idx)
                        row_text = ' '.join([str(v) if v is not None else '' for v in row_values])
                        if row_text.strip():
                            texts.append(row_text)
                result = '\n'.join(texts)
                if result.strip():
                    return result, None
            except:
                pass
        temp_dir = tempfile.mkdtemp(prefix='fc_conv_')
        try:
            for fmt in ['txt:Text', 'docx', 'xlsx']:
                converted = TextExtractor._convert_via_libreoffice(filepath, temp_dir, fmt)
                if converted and os.path.exists(converted):
                    if fmt == 'txt:Text' or fmt == 'txt':
                        text, _ = TextExtractor._extract_text_file(converted)
                    elif fmt == 'docx':
                        text, err = TextExtractor._extract_docx(converted)
                    else:
                        text, err = TextExtractor._extract_xlsx(converted)
                    os.remove(converted)
                    if text.strip():
                        return text, None
            return "", "LibreOffice 转换失败"
        except Exception as e:
            return "", f"转换异常: {e}"
        finally:
            try:
                shutil.rmtree(temp_dir, ignore_errors=True)
            except:
                pass


# ==================== 匹配引擎 ====================
class MatchEngine:
    def __init__(self, keywords, regex_patterns):
        self.keywords = keywords
        self.regex_patterns = regex_patterns

    def search(self, content):
        matches = []
        content_lower = content.lower()
        for keyword in self.keywords:
            idx = content_lower.find(keyword.lower())
            if idx >= 0:
                start = max(0, idx - 50)
                end = min(len(content), idx + len(keyword) + 50)
                context = content[start:end].replace('\n', ' ').replace('\r', ' ')
                matches.append(('关键词', keyword, context))
        for pattern, compiled_regex in self.regex_patterns:
            for m in compiled_regex.finditer(content):
                matched_text = m.group()
                start = max(0, m.start() - 50)
                end = min(len(content), m.end() + 50)
                context = content[start:end].replace('\n', ' ').replace('\r', ' ')
                matches.append(('正则表达式', pattern, context))
        return matches


# ==================== 文件检查器 ====================
class FileChecker:
    def __init__(self, keywords, regex_patterns, config=None):
        self.engine = MatchEngine(keywords, regex_patterns)
        self.config = config or {}
        self.max_file_size = self.config.get('max_file_size', MAX_FILE_SIZE_DEFAULT)
        self.scan_hidden = self.config.get('scan_hidden', False)
        self.scan_zip = self.config.get('scan_zip', True)
        self.enable_ocr = self.config.get('enable_ocr', True)
        self.results = []
        self.stats = {'total_files': 0, 'scanned_files': 0, 'skipped_files': 0,
                      'matched_files': 0, 'total_matches': 0, 'errors': []}
        self._stop_flag = False

    def stop(self):
        self._stop_flag = True

    def should_skip_dir(self, dirname):
        if not self.scan_hidden and dirname.startswith('.'):
            return True
        return dirname in SKIP_DIRS

    def should_skip_file(self, filepath, ext_lower):
        if ext_lower in SKIP_EXTENSIONS:
            return True, "跳过系统/二进制文件类型"
        try:
            file_size = os.path.getsize(filepath)
            if file_size > self.max_file_size:
                return True, f"文件过大 ({file_size // (1024*1024)}MB)"
            if file_size == 0:
                return True, "空文件"
        except OSError:
            return True, "无法获取文件大小"
        supported = (TEXT_EXTENSIONS | DOCX_EXTENSIONS | XLSX_EXTENSIONS |
                     PPTX_EXTENSIONS | DOC_EXTENSIONS | XLS_EXTENSIONS | PPT_EXTENSIONS |
                     PDF_EXTENSIONS | EML_EXTENSIONS | ZIP_EXTENSIONS |
                     WPS_TEXT_EXTENSIONS | WPS_SHEET_EXTENSIONS | WPS_SLIDE_EXTENSIONS |
                     IMAGE_EXTENSIONS | OFD_EXTENSIONS)
        if ext_lower not in supported:
            return True, "不支持的文件类型"
        return False, None

    def scan_file(self, filepath):
        if self._stop_flag:
            return 0
        ext_lower = os.path.splitext(filepath)[1].lower()
        if ext_lower in ZIP_EXTENSIONS and self.scan_zip:
            return self._scan_zip_file(filepath)
        text, error = TextExtractor.extract_text(filepath, ext_lower)
        if error:
            self.stats['errors'].append(f"{filepath}: {error}")
            return 0
        if not text or not text.strip():
            return 0
        matches = self.engine.search(text)
        for m_type, keyword, context in matches:
            self.results.append({
                'filepath': filepath,
                'filename': os.path.basename(filepath),
                'match_type': m_type,
                'keyword': keyword,
                'context': context
            })
        return len(matches)

    def _scan_zip_file(self, zip_path):
        total_matches = 0
        try:
            with zipfile.ZipFile(zip_path, 'r') as zf:
                temp_dir = tempfile.mkdtemp(prefix='fc_zip_')
                try:
                    for info in zf.infolist():
                        if self._stop_flag:
                            break
                        if info.is_dir() or info.flag_bits & 0x1:
                            continue
                        if info.file_size > self.max_file_size:
                            continue
                        try:
                            extracted_path = zf.extract(info, temp_dir)
                            total_matches += self.scan_file(extracted_path)
                        except Exception as e:
                            self.stats['errors'].append(f"{zip_path}::{info.filename}: {e}")
                        finally:
                            if 'extracted_path' in dir():
                                try:
                                    os.remove(extracted_path)
                                except:
                                    pass
                finally:
                    shutil.rmtree(temp_dir, ignore_errors=True)
        except (zipfile.BadZipFile, Exception) as e:
            self.stats['errors'].append(f"{zip_path}: 压缩文件读取失败 - {e}")
        return total_matches

    def scan_directory(self, scan_path, progress_callback=None):
        scan_path = os.path.abspath(scan_path)
        if not os.path.exists(scan_path):
            return
        if os.path.isfile(scan_path):
            self.stats['total_files'] = 1
            skip, reason = self.should_skip_file(scan_path,
                                                  os.path.splitext(scan_path)[1].lower())
            if not skip:
                self.stats['scanned_files'] = 1
                matches = self.scan_file(scan_path)
                if matches > 0:
                    self.stats['matched_files'] = 1
                    self.stats['total_matches'] += matches
            return

        file_list = []
        for root, dirs, files in os.walk(scan_path):
            dirs[:] = [d for d in dirs if not self.should_skip_dir(d)]
            for filename in files:
                filepath = os.path.join(root, filename)
                file_list.append(filepath)
        total = len(file_list)
        self.stats['total_files'] = total

        for idx, filepath in enumerate(file_list):
            if self._stop_flag:
                break
            # ★ 开始处理前立即更新进度（让用户看到当前正在处理哪个文件）
            if progress_callback:
                progress_callback(idx + 1, total, filepath)
            ext_lower = os.path.splitext(filepath)[1].lower()
            skip, reason = self.should_skip_file(filepath, ext_lower)
            if skip:
                self.stats['skipped_files'] += 1
            else:
                self.stats['scanned_files'] += 1
                try:
                    matches = self.scan_file(filepath)
                    if matches > 0:
                        self.stats['matched_files'] += 1
                        self.stats['total_matches'] += matches
                except Exception as e:
                    self.stats['errors'].append(f"{filepath}: {e}")


# ==================== 主 GUI 应用 ====================
class FileCheckerGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("文件内容字典检索工具 V5.1")
        self.root.geometry("1100x750")
        self.root.minsize(900, 600)
        self.scan_thread = None
        self.checker = None
        self.scanning = False
        self.style = ttk.Style()
        self.style.theme_use('clam')
        self._build_ui()
        self.root.update_idletasks()
        w = self.root.winfo_width()
        h = self.root.winfo_height()
        x = (self.root.winfo_screenwidth() // 2) - (w // 2)
        y = (self.root.winfo_screenheight() // 2) - (h // 2)
        self.root.geometry(f'+{x}+{y}')

    def _build_ui(self):
        main_frame = ttk.Frame(self.root, padding=10)
        main_frame.pack(fill=tk.BOTH, expand=True)

        # ===== 顶部：扫描设置 =====
        top_frame = ttk.LabelFrame(main_frame, text="扫描设置", padding=10)
        top_frame.pack(fill=tk.X, pady=(0, 10))

        row1 = ttk.Frame(top_frame)
        row1.pack(fill=tk.X, pady=2)
        ttk.Label(row1, text="扫描目录:").pack(side=tk.LEFT)
        self.dir_var = tk.StringVar()
        dir_entry = ttk.Entry(row1, textvariable=self.dir_var, width=60)
        dir_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=5)
        ttk.Button(row1, text="浏览...", command=self._browse_dir).pack(side=tk.LEFT)
        ttk.Button(row1, text="浏览文件...", command=self._browse_file).pack(side=tk.LEFT, padx=5)

        row2 = ttk.Frame(top_frame)
        row2.pack(fill=tk.X, pady=5)
        ttk.Label(row2, text="最大文件大小:").pack(side=tk.LEFT)
        self.size_var = tk.StringVar(value="50")
        ttk.Entry(row2, textvariable=self.size_var, width=6).pack(side=tk.LEFT, padx=5)
        ttk.Label(row2, text="MB").pack(side=tk.LEFT)
        self.scan_zip_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(row2, text="扫描压缩文件", variable=self.scan_zip_var).pack(side=tk.LEFT, padx=15)
        self.scan_hidden_var = tk.BooleanVar(value=False)
        ttk.Checkbutton(row2, text="扫描隐藏文件", variable=self.scan_hidden_var).pack(side=tk.LEFT)
        self.enable_ocr_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(row2, text="启用 OCR（图片/扫描件）", variable=self.enable_ocr_var).pack(side=tk.LEFT, padx=15)

        # ===== 中间：关键词字典 =====
        mid_frame = ttk.LabelFrame(main_frame, text="关键词字典（每行一个关键词，支持 regex: 前缀）", padding=10)
        mid_frame.pack(fill=tk.BOTH, expand=True, pady=(0, 10))

        dict_frame = ttk.Frame(mid_frame)
        dict_frame.pack(fill=tk.BOTH, expand=True)

        self.dict_text = tk.Text(dict_frame, wrap=tk.NONE, font=("Consolas", 11),
                                 undo=True, width=40, height=10)
        dict_scroll_y = ttk.Scrollbar(dict_frame, orient=tk.VERTICAL, command=self.dict_text.yview)
        dict_scroll_x = ttk.Scrollbar(dict_frame, orient=tk.HORIZONTAL, command=self.dict_text.xview)
        self.dict_text.configure(yscrollcommand=dict_scroll_y.set, xscrollcommand=dict_scroll_x.set)
        self.dict_text.grid(row=0, column=0, sticky='nsew')
        dict_scroll_y.grid(row=0, column=1, sticky='ns')
        dict_scroll_x.grid(row=1, column=0, sticky='ew')
        dict_frame.grid_rowconfigure(0, weight=1)
        dict_frame.grid_columnconfigure(0, weight=1)

        dict_btn_frame = ttk.Frame(mid_frame)
        dict_btn_frame.pack(fill=tk.X, pady=(5, 0))
        ttk.Button(dict_btn_frame, text="加载字典文件...", command=self._load_dict_file).pack(side=tk.LEFT)
        ttk.Button(dict_btn_frame, text="保存字典...", command=self._save_dict_file).pack(side=tk.LEFT, padx=5)
        ttk.Button(dict_btn_frame, text="清空", command=lambda: self.dict_text.delete('1.0', tk.END)).pack(side=tk.LEFT)
        ttk.Button(dict_btn_frame, text="加载示例", command=self._load_example).pack(side=tk.LEFT, padx=5)
        ttk.Label(dict_btn_frame, text="（# 开头的行视为注释，将被忽略）",
                  foreground="gray").pack(side=tk.RIGHT)

        # ===== 底部操作栏 =====
        ctrl_frame = ttk.Frame(main_frame)
        ctrl_frame.pack(fill=tk.X, pady=(0, 5))

        btn_frame = ttk.Frame(ctrl_frame)
        btn_frame.pack(side=tk.LEFT)
        self.start_btn = ttk.Button(btn_frame, text="▶ 开始扫描", command=self._start_scan, width=12)
        self.start_btn.pack(side=tk.LEFT, padx=2)
        self.stop_btn = ttk.Button(btn_frame, text="■ 停止", command=self._stop_scan, state=tk.DISABLED, width=8)
        self.stop_btn.pack(side=tk.LEFT, padx=2)
        self.export_btn = ttk.Button(btn_frame, text="📄 导出报告", command=self._export_report, state=tk.DISABLED, width=12)
        self.export_btn.pack(side=tk.LEFT, padx=2)
        self.export_paths_var = tk.BooleanVar(value=False)
        ttk.Checkbutton(btn_frame, text="同时导出路径列表", variable=self.export_paths_var).pack(side=tk.LEFT, padx=5)

        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(ctrl_frame, variable=self.progress_var, maximum=100)
        self.progress_bar.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=10)

        self.status_var = tk.StringVar(value="就绪")
        ttk.Label(ctrl_frame, textvariable=self.status_var, foreground="gray").pack(side=tk.RIGHT)

        # ===== 底部：结果列表 =====
        result_frame = ttk.LabelFrame(main_frame, text="扫描结果", padding=10)
        result_frame.pack(fill=tk.BOTH, expand=True)

        tree_frame = ttk.Frame(result_frame)
        tree_frame.pack(fill=tk.BOTH, expand=True)

        columns = ('#', '文件路径', '匹配类型', '关键词', '上下文')
        self.tree = ttk.Treeview(tree_frame, columns=columns, show='headings', selectmode='extended')
        self.tree.heading('#', text='#', command=lambda: self._sort_by_col(0))
        self.tree.heading('文件路径', text='文件路径', command=lambda: self._sort_by_col(1))
        self.tree.heading('匹配类型', text='类型')
        self.tree.heading('关键词', text='关键词')
        self.tree.heading('上下文', text='上下文')
        self.tree.column('#', width=40, minwidth=40)
        self.tree.column('文件路径', width=300, minwidth=100)
        self.tree.column('匹配类型', width=80, minwidth=60)
        self.tree.column('关键词', width=120, minwidth=80)
        self.tree.column('上下文', width=400, minwidth=150)

        tree_scroll_y = ttk.Scrollbar(tree_frame, orient=tk.VERTICAL, command=self.tree.yview)
        tree_scroll_x = ttk.Scrollbar(tree_frame, orient=tk.HORIZONTAL, command=self.tree.xview)
        self.tree.configure(yscrollcommand=tree_scroll_y.set, xscrollcommand=tree_scroll_x.set)
        self.tree.grid(row=0, column=0, sticky='nsew')
        tree_scroll_y.grid(row=0, column=1, sticky='ns')
        tree_scroll_x.grid(row=1, column=0, sticky='ew')
        tree_frame.grid_rowconfigure(0, weight=1)
        tree_frame.grid_columnconfigure(0, weight=1)
        self.tree.bind('<Double-1>', self._on_tree_double_click)

        stats_frame = ttk.Frame(result_frame)
        stats_frame.pack(fill=tk.X, pady=(5, 0))
        self.stats_var = tk.StringVar(value="扫描: 0 | 跳过: 0 | 命中文件: 0 | 命中次数: 0")
        ttk.Label(stats_frame, textvariable=self.stats_var, foreground="navy").pack(side=tk.LEFT)

        ttk.Label(main_frame, text="本工具为通用文件内容检索工具，仅供合法合规使用",
                  foreground="gray", font=("", 8)).pack(pady=(5, 0))

    # ==================== 事件处理 ====================
    def _browse_dir(self):
        path = filedialog.askdirectory(title="选择扫描目录")
        if path:
            self.dir_var.set(path)

    def _browse_file(self):
        path = filedialog.askopenfilename(title="选择单个文件")
        if path:
            self.dir_var.set(path)

    def _load_dict_file(self):
        path = filedialog.askopenfilename(title="打开字典文件",
                                          filetypes=[("文本文件", "*.txt"), ("所有文件", "*.*")])
        if path:
            try:
                with open(path, 'r', encoding='utf-8') as f:
                    content = f.read()
                self.dict_text.delete('1.0', tk.END)
                self.dict_text.insert('1.0', content)
            except Exception as e:
                messagebox.showerror("错误", f"加载字典文件失败:\n{e}")

    def _save_dict_file(self):
        path = filedialog.asksaveasfilename(title="保存字典文件",
                                            defaultextension=".txt",
                                            filetypes=[("文本文件", "*.txt")])
        if path:
            try:
                with open(path, 'w', encoding='utf-8') as f:
                    f.write(self.dict_text.get('1.0', tk.END).rstrip())
            except Exception as e:
                messagebox.showerror("错误", f"保存失败:\n{e}")

    def _load_example(self):
        example = """# --- 中文密级标志 ---
绝密
机密
秘密
机密★
绝密★

# --- 内部文件标识 ---
内部资料
内部文件
不得外传
严禁外泄

# --- 正则表达式示例 ---
regex:[1-9]\\d{5}(19|20)\\d{2}(0[1-9]|1[0-2])(0[1-9]|[12]\\d|3[01])\\d{3}[\\dXx]
regex:1(3[0-9]|4[5-9]|5[0-35-9]|6[2567]|7[0-8]|8[0-9]|9[0-35-9])\\d{8}
regex:[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\\.[a-zA-Z]{2,}"""
        self.dict_text.delete('1.0', tk.END)
        self.dict_text.insert('1.0', example)

    def _parse_dictionary(self):
        raw_text = self.dict_text.get('1.0', tk.END)
        keywords = []
        regex_patterns = []
        for line_num, line in enumerate(raw_text.split('\n'), 1):
            line = line.strip()
            if not line or line.startswith('#'):
                continue
            if line.startswith('regex:'):
                pattern = line[6:].strip()
                if pattern:
                    try:
                        compiled = re.compile(pattern, re.IGNORECASE)
                        regex_patterns.append((pattern, compiled))
                    except re.error as e:
                        messagebox.showwarning("正则表达式错误",
                                               f"第{line_num}行正则无效:\n{pattern}\n\n{e}")
                        return None, None
            else:
                keywords.append(line)
        return keywords, regex_patterns

    def _start_scan(self):
        if self.scanning:
            return
        scan_path = self.dir_var.get().strip()
        if not scan_path:
            messagebox.showwarning("警告", "请先选择扫描目录或文件")
            return
        if not os.path.exists(scan_path):
            messagebox.showerror("错误", f"路径不存在:\n{scan_path}")
            return
        keywords, regex_patterns = self._parse_dictionary()
        if keywords is None:
            return
        if not keywords and not regex_patterns:
            messagebox.showwarning("警告", "请至少输入一个关键词或正则表达式")
            return
        try:
            max_size_mb = int(self.size_var.get())
            if max_size_mb <= 0:
                raise ValueError
        except ValueError:
            messagebox.showwarning("警告", "最大文件大小必须为正整数")
            return
        for item in self.tree.get_children():
            self.tree.delete(item)
        config = {
            'max_file_size': max_size_mb * 1024 * 1024,
            'scan_hidden': self.scan_hidden_var.get(),
            'scan_zip': self.scan_zip_var.get(),
            'enable_ocr': self.enable_ocr_var.get(),
        }
        self.checker = FileChecker(keywords, regex_patterns, config)
        self.scanning = True
        self.start_btn.config(state=tk.DISABLED)
        self.stop_btn.config(state=tk.NORMAL)
        self.export_btn.config(state=tk.DISABLED)
        self.progress_var.set(0)
        self.status_var.set("正在扫描...")
        self.scan_thread = threading.Thread(target=self._scan_worker, args=(scan_path,), daemon=True)
        self.scan_thread.start()
        self._poll_thread()

    def _scan_worker(self, scan_path):
        def progress_callback(current, total, current_file):
            percent = (current / total * 100) if total > 0 else 0
            ext = os.path.splitext(current_file)[1].lower()
            self.root.after(0, self._update_progress, percent, current_file, ext)
        self.checker.scan_directory(scan_path, progress_callback)

    def _update_progress(self, percent, current_file, ext=''):
        self.progress_var.set(percent)
        base = os.path.basename(current_file)
        # OCR 文件给特殊提示
        if ext in ('.pdf', '.jpg', '.jpeg', '.png', '.bmp', '.tif', '.tiff') and self.enable_ocr_var.get():
            self.status_var.set(f"OCR 识别中: {base} （可能较慢）...")
        else:
            self.status_var.set(f"扫描中: {base}")

    def _poll_thread(self):
        if self.scan_thread is None:
            return
        if self.scan_thread.is_alive():
            self.root.after(200, self._poll_thread)
        else:
            self._scan_finished()

    def _scan_finished(self):
        self.scanning = False
        self.start_btn.config(state=tk.NORMAL)
        self.stop_btn.config(state=tk.DISABLED)
        if self.checker is None:
            return
        s = self.checker.stats
        self.status_var.set(f"完成 - 命中 {s['matched_files']} 个文件, {s['total_matches']} 次")
        self.stats_var.set(f"扫描: {s['scanned_files']} | 跳过: {s['skipped_files']} | "
                           f"命中文件: {s['matched_files']} | 命中次数: {s['total_matches']}")
        for i, r in enumerate(self.checker.results, 1):
            self.tree.insert('', tk.END, values=(i, r['filepath'], r['match_type'],
                                                  r['keyword'], r['context']))
        if self.checker.results:
            self.export_btn.config(state=tk.NORMAL)
        else:
            messagebox.showinfo("提示", "未找到匹配项")
        if s['errors']:
            self.status_var.set(self.status_var.get() + f" | 错误: {len(s['errors'])}")

    def _stop_scan(self):
        if self.checker:
            self.checker.stop()
        self.status_var.set("正在停止...")
        self.stop_btn.config(state=tk.DISABLED)

    def _export_report(self):
        if not self.checker or not self.checker.results:
            return
        filetypes = [("HTML 报告", "*.html"), ("CSV 文件", "*.csv"), ("文本文件", "*.txt")]
        path = filedialog.asksaveasfilename(title="导出报告", defaultextension=".html", filetypes=filetypes)
        if not path:
            return
        ext = os.path.splitext(path)[1].lower()
        try:
            if ext == '.csv':
                self._export_csv(path)
            elif ext == '.txt':
                self._export_text(path)
            else:
                self._export_html(path)
            if self.export_paths_var.get():
                paths_file = os.path.splitext(path)[0] + ".paths.txt"
                self._export_path_list(paths_file)
                messagebox.showinfo("成功", f"报告已保存到:\n{path}\n\n路径列表已保存到:\n{paths_file}")
            else:
                messagebox.showinfo("成功", f"报告已保存到:\n{path}")
        except Exception as e:
            messagebox.showerror("错误", f"导出失败:\n{e}")

    def _export_html(self, path):
        results = self.checker.results
        stats = self.checker.stats
        file_groups = {}
        for r in results:
            fp = r['filepath']
            if fp not in file_groups:
                file_groups[fp] = []
            file_groups[fp].append(r)
        with open(path, 'w', encoding='utf-8') as f:
            f.write(f'''<!DOCTYPE html><html lang="zh-CN"><head><meta charset="UTF-8">
<title>文件内容检查报告</title>
<style>body{{font-family:"Microsoft YaHei","SimSun",sans-serif;margin:20px;background:#f5f5f5}}
.container{{max-width:1200px;margin:0 auto;background:#fff;padding:30px;border-radius:8px;box-shadow:0 2px 10px rgba(0,0,0,0.1)}}
h1{{color:#333;border-bottom:2px solid #1a73e8;padding-bottom:10px}}
.summary{{background:#e8f0fe;padding:15px;border-radius:5px;margin:20px 0}}
.summary table{{border-collapse:collapse;width:100%}}.summary td{{padding:5px 10px}}
.file-group{{margin:15px 0;border:1px solid #ddd;border-radius:5px;overflow:hidden}}
.file-header{{background:#f8f9fa;padding:10px 15px;font-weight:bold}}
.match-table{{width:100%;border-collapse:collapse}}
.match-table th{{background:#1a73e8;color:#fff;padding:8px;text-align:left}}
.match-table td{{padding:8px;border-bottom:1px solid #eee}}
.match-table tr:hover{{background:#f1f3f4}}
.keyword{{color:#d93025;font-weight:bold}}
.context{{font-size:0.9em;color:#555;max-width:500px;word-break:break-all}}
.footer{{text-align:center;color:#999;margin-top:30px;font-size:0.85em}}</style></head>
<body><div class="container"><h1>文件内容字典检索报告</h1>
<div class="summary"><h3>检查摘要</h3><table>
<tr><td><strong>扫描目录:</strong></td><td>{escape(self.dir_var.get())}</td></tr>
<tr><td><strong>扫描时间:</strong></td><td>{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</td></tr>
<tr><td><strong>扫描文件:</strong></td><td>{stats['scanned_files']}</td></tr>
<tr><td><strong>跳过文件:</strong></td><td>{stats['skipped_files']}</td></tr>
<tr><td><strong>命中文件:</strong></td><td style="color:#d93025;font-weight:bold;">{stats['matched_files']}</td></tr>
<tr><td><strong>命中次数:</strong></td><td style="color:#d93025;font-weight:bold;">{stats['total_matches']}</td></tr>
</table></div><h2>匹配详情</h2>''')
            for filepath, matches in file_groups.items():
                f.write(f'<div class="file-group"><div class="file-header">{escape(filepath)} （{len(matches)} 处）</div><table class="match-table"><tr><th>#</th><th>类型</th><th>关键词</th><th>上下文</th></tr>')
                for i, m in enumerate(matches, 1):
                    f.write(f'<tr><td>{i}</td><td>{escape(m["match_type"])}</td><td class="keyword">{escape(m["keyword"])}</td><td class="context">{escape(m["context"])}</td></tr>')
                f.write('</table></div>')
            f.write('<div class="footer"><p>本工具为通用文件内容检索工具，仅供合法合规使用</p></div></div></body></html>')

    def _export_csv(self, path):
        with open(path, 'w', newline='', encoding='utf-8-sig') as f:
            writer = csv.writer(f)
            writer.writerow(['序号', '文件路径', '文件名', '匹配类型', '匹配关键词', '上下文'])
            for i, r in enumerate(self.checker.results, 1):
                writer.writerow([i, r['filepath'], r['filename'],
                                 r['match_type'], r['keyword'], r['context']])

    def _export_text(self, path):
        lines = ['=' * 70, '  文件内容字典检索报告', '=' * 70,
                 f'扫描目录: {self.dir_var.get()}',
                 f'扫描时间: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}',
                 f'扫描文件: {self.checker.stats["scanned_files"]} 个',
                 f'命中文件: {self.checker.stats["matched_files"]} 个',
                 f'命中次数: {self.checker.stats["total_matches"]} 次', '=' * 70, '']
        file_groups = {}
        for r in self.checker.results:
            fp = r['filepath']
            if fp not in file_groups:
                file_groups[fp] = []
            file_groups[fp].append(r)
        for filepath, matches in file_groups.items():
            lines.append(f'[文件] {filepath} ({len(matches)}处匹配)')
            for i, m in enumerate(matches, 1):
                lines.append(f'  {i}. [{m["match_type"]}] {m["keyword"]}')
                lines.append(f'     上下文: ...{m["context"]}...')
            lines.append('')
        lines.append('=' * 70)
        with open(path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(lines))

    def _export_path_list(self, filepath):
        file_set = set()
        for r in self.checker.results:
            file_set.add(r['filepath'])
        with open(filepath, 'w', encoding='utf-8') as f:
            for fp in sorted(file_set):
                f.write(fp + '\n')

    def _sort_by_col(self, col):
        rows = [(self.tree.set(item, '#'), item) for item in self.tree.get_children('')]
        try:
            rows.sort(key=lambda x: int(x[0]))
        except ValueError:
            rows.sort(key=lambda x: x[0])
        for idx, (_, item) in enumerate(rows):
            self.tree.move(item, '', idx)
            self.tree.set(item, '#', str(idx + 1))

    def _on_tree_double_click(self, event):
        item = self.tree.focus()
        if not item:
            return
        values = self.tree.item(item, 'values')
        if values and len(values) >= 2:
            filepath = values[1]
            if os.path.exists(filepath):
                if platform.system() == 'Windows':
                    os.system(f'explorer /select,"{filepath}"')
                elif platform.system() == 'Linux':
                    dirname = os.path.dirname(filepath)
                    subprocess.Popen(['xdg-open', dirname])
                elif platform.system() == 'Darwin':
                    subprocess.Popen(['open', '-R', filepath])


# ==================== 启动入口 ====================
def main():
    root = tk.Tk()
    app = FileCheckerGUI(root)
    root.mainloop()


if __name__ == '__main__':
    main()